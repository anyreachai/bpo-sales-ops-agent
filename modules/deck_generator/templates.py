"""shadcn-inspired design system for pitch deck generation.

Zinc neutral base with brand accent theming. Consistent grid system,
rounded shapes, and professional typography.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Default palette — zinc neutral base
# ---------------------------------------------------------------------------

DEFAULT_PALETTE = {
    "dark_bg": "#09090B",
    "light_bg": "#FAFAFA",
    "muted_bg": "#F4F4F5",
    "primary_accent": "#5B5FC7",
    "secondary_accent": "#818CF8",
    "neutral": ["#71717A", "#A1A1AA", "#D4D4D8"],
}

# ---------------------------------------------------------------------------
# Typography
# ---------------------------------------------------------------------------

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"

TITLE_SIZE = Pt(36)
SUBTITLE_SIZE = Pt(20)
HEADING_SIZE = Pt(28)
SUBHEADING_SIZE = Pt(18)
BODY_SIZE = Pt(14)
BULLET_SIZE = Pt(13)
SMALL_SIZE = Pt(11)
STAT_NUMBER_SIZE = Pt(44)
STAT_LABEL_SIZE = Pt(12)
FOOTER_SIZE = Pt(8)
CAPTION_SIZE = Pt(9)

# ---------------------------------------------------------------------------
# Layout constants (16:9 widescreen)
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.7)
CONTENT_LEFT = MARGIN
CONTENT_TOP = Inches(1.4)
CONTENT_WIDTH = Inches(11.933)
CONTENT_HEIGHT = Inches(5.0)

ACCENT_BAR_HEIGHT = Inches(0.05)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _hex(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_bg(slide, hex_color: str) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex(hex_color)


def _accent_bar(slide, palette: dict) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_W, ACCENT_BAR_HEIGHT,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(palette["primary_accent"])
    shape.line.fill.background()


def _bottom_bar(slide, palette: dict) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), SLIDE_H - Inches(0.04),
        SLIDE_W, Inches(0.04),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(palette["primary_accent"])
    shape.line.fill.background()


def _text_box(
    slide, left, top, width, height,
    text: str,
    font_name: str = BODY_FONT,
    font_size=BODY_SIZE,
    font_color: str = "#3F3F46",
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,
) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = _hex(font_color)
    p.font.bold = bold
    p.alignment = alignment


def _bullets(
    slide, left, top, width, height,
    items: list[str],
    font_color: str = "#3F3F46",
    font_size=BULLET_SIZE,
) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = font_size
        p.font.color.rgb = _hex(font_color)
        p.space_after = Pt(10)
        p.level = 0
        pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        for existing in pPr.findall(qn("a:buChar")):
            pPr.remove(existing)
        for existing in pPr.findall(qn("a:buNone")):
            pPr.remove(existing)
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
        pPr.append(buChar)


def _rounded_card(slide, left, top, width, height, fill_color: str, border_color: str | None = None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(fill_color)
    if border_color:
        shape.line.color.rgb = _hex(border_color)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _footer(slide, palette: dict, company: str = "", slide_num: int = 0, total: int = 0) -> None:
    neutral = palette.get("neutral", DEFAULT_PALETTE["neutral"])
    color = neutral[1] if isinstance(neutral, list) and len(neutral) > 1 else "#A1A1AA"

    footer_text = "Confidential | Anyreach.ai"
    if company:
        footer_text = f"{company} | Confidential | Anyreach.ai"

    _text_box(
        slide,
        left=MARGIN,
        top=Inches(7.05),
        width=Inches(6.0),
        height=Inches(0.35),
        text=footer_text,
        font_size=FOOTER_SIZE,
        font_color=color,
        alignment=PP_ALIGN.LEFT,
    )

    if slide_num and total:
        _text_box(
            slide,
            left=Inches(11.0),
            top=Inches(7.05),
            width=Inches(2.0),
            height=Inches(0.35),
            text=f"{slide_num} / {total}",
            font_size=FOOTER_SIZE,
            font_color=color,
            alignment=PP_ALIGN.RIGHT,
        )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def add_title_slide(prs, slide_data: dict, palette: dict, logo_path=None, **kw) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, palette["dark_bg"])

    # Full-width accent bar
    _accent_bar(slide, palette)
    _bottom_bar(slide, palette)

    # Title
    _text_box(
        slide,
        left=Inches(1.0), top=Inches(2.2),
        width=Inches(11.0), height=Inches(1.4),
        text=slide_data.get("title", ""),
        font_name=TITLE_FONT,
        font_size=TITLE_SIZE,
        font_color="#FAFAFA",
        bold=True,
    )

    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        _text_box(
            slide,
            left=Inches(1.0), top=Inches(3.7),
            width=Inches(11.0), height=Inches(0.8),
            text=subtitle,
            font_size=SUBTITLE_SIZE,
            font_color=palette.get("secondary_accent", "#818CF8"),
        )

    prepared_for = slide_data.get("prepared_for", "")
    if prepared_for:
        _text_box(
            slide,
            left=Inches(1.0), top=Inches(5.0),
            width=Inches(11.0), height=Inches(0.5),
            text=prepared_for,
            font_size=SMALL_SIZE,
            font_color="#A1A1AA",
        )

    if logo_path:
        try:
            slide.shapes.add_picture(logo_path, Inches(10.0), Inches(5.8), Inches(2.5), Inches(1.0))
        except Exception:
            pass


def add_content_slide(prs, slide_data: dict, palette: dict, **kw) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, palette["light_bg"])
    _accent_bar(slide, palette)

    heading = slide_data.get("heading", "")
    _text_box(
        slide,
        left=MARGIN, top=Inches(0.5),
        width=CONTENT_WIDTH, height=Inches(0.7),
        text=heading,
        font_name=TITLE_FONT,
        font_size=HEADING_SIZE,
        font_color=palette["dark_bg"],
        bold=True,
    )

    # Thin divider under heading
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN, Inches(1.15),
        Inches(3.0), Inches(0.025),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = _hex(palette["primary_accent"])
    divider.line.fill.background()

    subheading = slide_data.get("subheading", "")
    body_top = CONTENT_TOP
    if subheading:
        _text_box(
            slide,
            left=MARGIN, top=Inches(1.2),
            width=CONTENT_WIDTH, height=Inches(0.4),
            text=subheading,
            font_size=SUBHEADING_SIZE,
            font_color=palette.get("primary_accent", "#5B5FC7"),
        )
        body_top = Inches(1.75)

    bullets_list = slide_data.get("bullets", [])
    if bullets_list:
        _bullets(
            slide,
            left=Inches(0.9), top=body_top,
            width=Inches(11.0), height=CONTENT_HEIGHT,
            items=bullets_list,
            font_color="#3F3F46",
        )

    body = slide_data.get("body", "")
    if body and not bullets_list:
        _text_box(
            slide,
            left=MARGIN, top=body_top,
            width=CONTENT_WIDTH, height=CONTENT_HEIGHT,
            text=body,
            font_color="#3F3F46",
        )


def add_stats_slide(prs, slide_data: dict, palette: dict, **kw) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, palette["light_bg"])
    _accent_bar(slide, palette)

    heading = slide_data.get("heading", "")
    _text_box(
        slide,
        left=MARGIN, top=Inches(0.5),
        width=CONTENT_WIDTH, height=Inches(0.7),
        text=heading,
        font_name=TITLE_FONT,
        font_size=HEADING_SIZE,
        font_color=palette["dark_bg"],
        bold=True,
    )

    stats = slide_data.get("stats", [])
    if not stats:
        return

    # 2x2 grid of rounded cards
    count = min(len(stats), 4)
    card_w = Inches(5.0)
    card_h = Inches(2.2)
    gap_x = Inches(0.8)
    gap_y = Inches(0.5)

    if count <= 3:
        # Single row
        total_w = count * card_w + (count - 1) * gap_x
        start_x = (SLIDE_W - total_w) / 2
        positions = [(start_x + i * (card_w + gap_x), Inches(2.4)) for i in range(count)]
    else:
        # 2x2 grid
        start_x = (SLIDE_W - 2 * card_w - gap_x) / 2
        positions = [
            (start_x, Inches(1.8)),
            (start_x + card_w + gap_x, Inches(1.8)),
            (start_x, Inches(1.8) + card_h + gap_y),
            (start_x + card_w + gap_x, Inches(1.8) + card_h + gap_y),
        ]

    for i, stat in enumerate(stats[:4]):
        left, top = positions[i]

        # Card background
        _rounded_card(slide, left, top, card_w, card_h, palette.get("muted_bg", "#F4F4F5"), "#D4D4D8")

        # Accent stripe on card
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + Inches(0.15), top + Inches(0.15),
            Inches(0.06), card_h - Inches(0.3),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = _hex(palette["primary_accent"])
        stripe.line.fill.background()

        # Stat value
        val = stat.get("value", stat.get("number", ""))
        _text_box(
            slide,
            left=left + Inches(0.4), top=top + Inches(0.35),
            width=card_w - Inches(0.6), height=Inches(1.0),
            text=str(val),
            font_name=TITLE_FONT,
            font_size=STAT_NUMBER_SIZE,
            font_color=palette["primary_accent"],
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Stat label
        _text_box(
            slide,
            left=left + Inches(0.4), top=top + Inches(1.3),
            width=card_w - Inches(0.6), height=Inches(0.7),
            text=stat.get("label", ""),
            font_size=STAT_LABEL_SIZE,
            font_color="#71717A",
            alignment=PP_ALIGN.LEFT,
        )


def add_comparison_slide(prs, slide_data: dict, palette: dict, **kw) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, palette["light_bg"])
    _accent_bar(slide, palette)

    heading = slide_data.get("heading", "")
    _text_box(
        slide,
        left=MARGIN, top=Inches(0.5),
        width=CONTENT_WIDTH, height=Inches(0.7),
        text=heading,
        font_name=TITLE_FONT,
        font_size=HEADING_SIZE,
        font_color=palette["dark_bg"],
        bold=True,
    )

    col_width = Inches(5.5)
    left_x = Inches(0.7)
    right_x = Inches(7.1)
    col_top = Inches(1.5)
    col_height = Inches(5.2)

    # Left panel — muted background
    _rounded_card(slide, left_x - Inches(0.1), col_top - Inches(0.1),
                  col_width + Inches(0.2), col_height + Inches(0.2),
                  "#F4F4F5", "#E4E4E7")

    # Right panel — accent tinted
    _rounded_card(slide, right_x - Inches(0.1), col_top - Inches(0.1),
                  col_width + Inches(0.2), col_height + Inches(0.2),
                  "#EEF2FF", "#C7D2FE")

    # Left column
    left_data = slide_data.get("left", {})
    _text_box(
        slide,
        left=left_x, top=col_top,
        width=col_width, height=Inches(0.5),
        text=left_data.get("title", "Current State"),
        font_size=SUBHEADING_SIZE,
        font_color="#71717A",
        bold=True,
    )
    left_bullets = left_data.get("bullets", [])
    if left_bullets:
        _bullets(
            slide,
            left=left_x, top=col_top + Inches(0.65),
            width=col_width, height=col_height - Inches(0.65),
            items=left_bullets,
            font_color="#52525B",
        )

    # Right column
    right_data = slide_data.get("right", {})
    _text_box(
        slide,
        left=right_x, top=col_top,
        width=col_width, height=Inches(0.5),
        text=right_data.get("title", "With Anyreach"),
        font_size=SUBHEADING_SIZE,
        font_color=palette["primary_accent"],
        bold=True,
    )
    right_bullets = right_data.get("bullets", [])
    if right_bullets:
        _bullets(
            slide,
            left=right_x, top=col_top + Inches(0.65),
            width=col_width, height=col_height - Inches(0.65),
            items=right_bullets,
            font_color="#27272A",
        )

    # Vertical divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.6), col_top,
        Inches(0.02), col_height,
    )
    neutral = palette.get("neutral", DEFAULT_PALETTE["neutral"])
    divider.fill.solid()
    divider.fill.fore_color.rgb = _hex(neutral[2] if isinstance(neutral, list) and len(neutral) > 2 else "#D4D4D8")
    divider.line.fill.background()


def add_quote_slide(prs, slide_data: dict, palette: dict, **kw) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, palette.get("muted_bg", "#F4F4F5"))
    _accent_bar(slide, palette)

    # Large quote mark
    _text_box(
        slide,
        left=Inches(1.5), top=Inches(1.5),
        width=Inches(1.0), height=Inches(1.0),
        text="\u201C",
        font_size=Pt(72),
        font_color=palette["primary_accent"],
        bold=True,
    )

    quote = slide_data.get("quote", slide_data.get("body", ""))
    _text_box(
        slide,
        left=Inches(2.0), top=Inches(2.2),
        width=Inches(9.0), height=Inches(2.5),
        text=quote,
        font_name=TITLE_FONT,
        font_size=Pt(22),
        font_color="#27272A",
        alignment=PP_ALIGN.LEFT,
    )

    attribution = slide_data.get("attribution", "")
    if attribution:
        _text_box(
            slide,
            left=Inches(2.0), top=Inches(5.0),
            width=Inches(9.0), height=Inches(0.5),
            text=f"— {attribution}",
            font_size=SMALL_SIZE,
            font_color="#71717A",
        )


def add_cta_slide(prs, slide_data: dict, palette: dict, logo_path=None, **kw) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, palette["dark_bg"])
    _accent_bar(slide, palette)
    _bottom_bar(slide, palette)

    heading = slide_data.get("heading", "Next Steps")
    _text_box(
        slide,
        left=Inches(1.5), top=Inches(1.5),
        width=Inches(10.0), height=Inches(1.0),
        text=heading,
        font_name=TITLE_FONT,
        font_size=Pt(32),
        font_color="#FAFAFA",
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    bullets_list = slide_data.get("bullets", [])
    if bullets_list:
        _bullets(
            slide,
            left=Inches(2.5), top=Inches(2.8),
            width=Inches(8.0), height=Inches(3.0),
            items=bullets_list,
            font_color="#D4D4D8",
            font_size=BODY_SIZE,
        )

    body = slide_data.get("body", "")
    if body and not bullets_list:
        _text_box(
            slide,
            left=Inches(2.5), top=Inches(2.8),
            width=Inches(8.0), height=Inches(3.0),
            text=body,
            font_color="#D4D4D8",
            alignment=PP_ALIGN.CENTER,
        )

    contact = slide_data.get("contact", "")
    if contact:
        _text_box(
            slide,
            left=Inches(1.5), top=Inches(5.8),
            width=Inches(10.0), height=Inches(0.5),
            text=contact,
            font_size=SUBTITLE_SIZE,
            font_color=palette.get("secondary_accent", "#818CF8"),
            alignment=PP_ALIGN.CENTER,
        )

    if logo_path:
        try:
            slide.shapes.add_picture(logo_path, Inches(5.5), Inches(6.3), Inches(2.0), Inches(0.8))
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Slide dispatcher
# ---------------------------------------------------------------------------

SLIDE_BUILDERS = {
    "title": add_title_slide,
    "content": add_content_slide,
    "stats_grid": add_stats_slide,
    "comparison": add_comparison_slide,
    "quote": add_quote_slide,
    "cta": add_cta_slide,
}


def create_presentation(
    slides_data: list[dict],
    palette: dict,
    company_name: str,
    bpo_name: str,
    logo_path: str | None = None,
    fonts: dict | None = None,
) -> Presentation:
    """Build the full .pptx from structured slide data and a color palette."""
    global TITLE_FONT, BODY_FONT
    if fonts:
        TITLE_FONT = fonts.get("heading", TITLE_FONT)
        BODY_FONT = fonts.get("body", BODY_FONT)

    # Ensure palette has all keys
    for key, default in DEFAULT_PALETTE.items():
        if key not in palette:
            palette[key] = default

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = len(slides_data)
    for idx, slide_data in enumerate(slides_data):
        slide_type = slide_data.get("type", "content")
        builder = SLIDE_BUILDERS.get(slide_type, add_content_slide)
        builder(prs, slide_data, palette, logo_path=logo_path)

        # Add footer to every slide
        slide = prs.slides[idx]
        _footer(slide, palette, company=company_name, slide_num=idx + 1, total=total)

    return prs
